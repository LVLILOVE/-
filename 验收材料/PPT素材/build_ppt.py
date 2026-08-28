# ============================================================
# 代码段功能：生成项目验收演示 PPT（python-pptx）
# - 结构：封面 → 项目背景 → 开发目的 → WorkBuddy 需求处理（截图×2）
#         → 项目效果（截图×3）→ 项目心得 → 结束页
# - 暖咖主题配色，16:9 宽屏
# ============================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- 主题色（猫屿暖咖系）----
BROWN = RGBColor(0x6E, 0x4F, 0x3A)     # 深咖（标题/强调）
ORANGE = RGBColor(0xC4, 0x6A, 0x33)    # 暖橙（主色）
CREAM = RGBColor(0xFA, 0xF6, 0xF0)     # 奶油底
INK = RGBColor(0x3D, 0x34, 0x2C)       # 正文深棕
SOFT = RGBColor(0x8A, 0x7D, 0x70)      # 次级文字
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ASSETS = r'C:\Users\Administrator\WorkBuddy\2026-08-26-14-15-18\-\验收材料\PPT素材'
OUT = r'C:\Users\Administrator\WorkBuddy\2026-08-26-14-15-18\-\验收材料\演示PPT'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_bg(slide, color):
    """整页背景色"""
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False

def textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, line_spacing=1.4):
    """通用文本框"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.bold = bold; p.alignment = align
        p.line_spacing = line_spacing
    return tb

def title_bar(slide, title, en):
    """内容页顶部标题条"""
    add_bg(slide, WHITE)
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    textbox(slide, Inches(0.6), Inches(0.35), Inches(11), Inches(0.7), title, 28, BROWN, True)
    textbox(slide, Inches(0.62), Inches(1.0), Inches(11), Inches(0.4), en, 12, SOFT)

def add_img(slide, path, x, y, w=None, h=None, border=True):
    """图片占位（等比缩放，可选边框）"""
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    if border:
        pic.line.color.rgb = RGBColor(0xE0, 0xD6, 0xC8); pic.line.width = Pt(1.5)
    return pic

# ================= 1. 封面 =================
s = prs.slides.add_slide(BLANK)
add_bg(s, BROWN)
textbox(s, Inches(2), Inches(1.2), Inches(9.3), Inches(1.0), '猫屿 CAT ISLE', 54, WHITE, True, PP_ALIGN.CENTER)
textbox(s, Inches(2), Inches(2.35), Inches(9.3), Inches(0.6), '猫 咖 官 网 · 项 目 验 收 演 示', 20, RGBColor(0xE8, 0xD8, 0xC2), False, PP_ALIGN.CENTER)
textbox(s, Inches(2), Inches(3.6), Inches(9.3), Inches(0.5), '在猫的节奏里，慢下来', 16, RGBColor(0xC9, 0xB5, 0x9C), False, PP_ALIGN.CENTER)
line = s.shapes.add_shape(1, Inches(4.67), Inches(4.4), Inches(4), Pt(1.2))
line.fill.solid(); line.fill.fore_color.rgb = ORANGE; line.line.fill.background()
textbox(s, Inches(2), Inches(5.4), Inches(9.3), Inches(0.5), '开发工具：WorkBuddy AI 助手', 15, RGBColor(0xD8, 0xC8, 0xB2), False, PP_ALIGN.CENTER)
textbox(s, Inches(2), Inches(6.1), Inches(9.3), Inches(0.5), '技术栈：FastAPI · React · SQLite · 2026-08', 13, RGBColor(0xB8, 0xA6, 0x8E), False, PP_ALIGN.CENTER)

# ================= 2. 项目背景 =================
s = prs.slides.add_slide(BLANK)
title_bar(s, '项目背景', 'PROJECT BACKGROUND')
bg_items = [
    ('01', '宠物消费升级', '越来越多年轻人把宠物当作家人，城市猫咖成为放松、社交的热门去处，但多数猫咖线上体验薄弱，仅靠朋友圈口口相传。'),
    ('02', '新客沟通成本高', '预约规则、押金流程、互动须知等问题反复被新客询问，缺少一个集中展示与自助答疑的线上窗口。'),
    ('03', '门店数字化缺失', '预约靠接龙、领养靠登记，数据散落，店长无法统一查看预约、押金与领养进度。'),
]
y = Inches(1.7)
for no, t, d in bg_items:
    num = s.shapes.add_shape(5, Inches(0.9), y + Inches(0.05), Inches(0.75), Inches(0.75))
    num.fill.solid(); num.fill.fore_color.rgb = ORANGE; num.line.fill.background()
    num.text_frame.text = no; num.text_frame.paragraphs[0].font.size = Pt(24)
    num.text_frame.paragraphs[0].font.color.rgb = WHITE; num.text_frame.paragraphs[0].font.bold = True
    num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    textbox(s, Inches(2.0), y, Inches(10.2), Inches(0.5), t, 20, BROWN, True)
    textbox(s, Inches(2.0), y + Inches(0.55), Inches(10.2), Inches(0.9), d, 14, INK)
    y += Inches(1.75)

# ================= 3. 开发目的 =================
s = prs.slides.add_slide(BLANK)
title_bar(s, '开发目的', 'DEVELOPMENT GOALS')
goals = [
    ('打造品牌线上门户', '以「猫屿」品牌为核心，统一呈现门店信息、猫咪成员、餐单与品牌故事，建立治愈系品牌认知。'),
    ('在线预约 + 押金闭环', '前台免注册预约，押金线下转账、店长后台核验，形成「预约-核验-到店」完整闭环，替代人工接龙。'),
    ('降低新客沟通成本', '店长解答问答区 + 猫屿小助手 + 今日店长，集中解答新客疑惑，减少重复人工咨询。'),
    ('领养流程数字化', '在线提交领养申请，店长四级审核 + 回访记录留痕，让领养更规范、可追溯。'),
    ('验收交付完整', '输出 PRD、技术文档、产品原型、数据库脚本、演示 PPT 全套验收材料。'),
]
y = Inches(1.6)
for i, (t, d) in enumerate(goals, 1):
    textbox(s, Inches(0.9), y, Inches(1.0), Inches(0.5), f'0{i}', 22, ORANGE, True)
    textbox(s, Inches(1.9), y, Inches(10.6), Inches(0.5), t, 17, BROWN, True)
    textbox(s, Inches(1.9), y + Inches(0.45), Inches(10.6), Inches(0.6), d, 13, INK)
    y += Inches(1.1)

# ================= 4. WorkBuddy 需求处理 =================
s = prs.slides.add_slide(BLANK)
title_bar(s, 'WorkBuddy 需求处理', 'AI-ASSISTED DEVELOPMENT')
textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
        '从需求提出到代码交付，全程由 WorkBuddy 智能助手处理：需求理解 → 方案设计 → 分阶段开发 → 测试验收。', 15, INK)
add_img(s, ASSETS + r'\04-需求处理记录.png', Inches(0.9), Inches(2.2), w=Inches(11.5), h=Inches(4.4))

s = prs.slides.add_slide(BLANK)
title_bar(s, 'WorkBuddy 需求处理（续）', 'AI-ASSISTED DEVELOPMENT')
add_img(s, ASSETS + r'\05-需求处理特写.png', Inches(1.7), Inches(1.8), w=Inches(10), h=Inches(4.6))

# ================= 5. 项目效果截图 =================
s = prs.slides.add_slide(BLANK)
title_bar(s, '项目效果（前台）', 'PROJECT SHOWCASE')
add_img(s, ASSETS + r'\01-首页.png', Inches(0.7), Inches(1.7), w=Inches(12), h=Inches(4.5))
textbox(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.4), '首页：Hero 品牌区 · 猫咪故事 · 猫咪预览 · 餐单预览 · 门店信息', 12, SOFT, False, PP_ALIGN.CENTER)

s = prs.slides.add_slide(BLANK)
title_bar(s, '项目效果（猫咪页）', 'PROJECT SHOWCASE')
add_img(s, ASSETS + r'\02-猫咪页.png', Inches(0.7), Inches(1.7), w=Inches(12), h=Inches(4.5))
textbox(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.4), '猫咪列表页：今日店长展示位 · 猫咪档案卡片 · 可领养筛选', 12, SOFT, False, PP_ALIGN.CENTER)

s = prs.slides.add_slide(BLANK)
title_bar(s, '项目效果（后台管理）', 'PROJECT SHOWCASE')
add_img(s, ASSETS + r'\03-后台管理.png', Inches(0.7), Inches(1.7), w=Inches(12), h=Inches(4.5))
textbox(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.4), '后台仪表盘：预约核验 · 押金管理 · 领养审核 · 门店配置 · 店长解答', 12, SOFT, False, PP_ALIGN.CENTER)

# ================= 6. 项目心得 =================
s = prs.slides.add_slide(BLANK)
title_bar(s, '项目心得', 'REFLECTIONS')
notes = [
    ('需求明确是成功的一半', '通过 PRD 与 UIUX 规范先行，把「猫咖预约、押金、领养」等复杂业务拆解成可执行的里程碑（M1-M4），开发过程几乎没有返工。'),
    ('AI 辅助开发 ≠ 全自动', 'WorkBuddy 大幅提升效率（建表、接口、页面、测试），但业务规则（押金核验、周一店休、领养审核）需要人工把关，代码注释与中文解释让审核更轻松。'),
    ('测试驱动更安心', '28 项 pytest 用例覆盖预约、押金、领养、问答等核心链路，端到端验证后交付，质量更有保障。'),
    ('部署思维前置', '单端口部署、内网穿透、本地图片存储等方案在设计阶段就确定，联调上线顺畅，验收材料水到渠成。'),
]
y = Inches(1.6)
for t, d in notes:
    textbox(s, Inches(0.9), y, Inches(11.5), Inches(0.5), '▪ ' + t, 17, BROWN, True)
    textbox(s, Inches(0.9), y + Inches(0.5), Inches(11.5), Inches(0.85), d, 13, INK)
    y += Inches(1.35)

# ================= 7. 结束页 =================
s = prs.slides.add_slide(BLANK)
add_bg(s, BROWN)
textbox(s, Inches(2), Inches(2.4), Inches(9.3), Inches(1.0), '感谢聆听', 44, WHITE, True, PP_ALIGN.CENTER)
textbox(s, Inches(2), Inches(3.8), Inches(9.3), Inches(0.6), '欢迎体验 · 敬请指正', 18, RGBColor(0xE8, 0xD8, 0xC2), False, PP_ALIGN.CENTER)
textbox(s, Inches(2), Inches(5.4), Inches(9.3), Inches(0.5), '猫屿 CAT ISLE · WorkBuddy 验收演示', 14, RGBColor(0xC9, 0xB5, 0x9C), False, PP_ALIGN.CENTER)

import os
os.makedirs(OUT, exist_ok=True)
out_path = os.path.join(OUT, '猫屿官网验收演示.pptx')
prs.save(out_path)
print('PPT 已生成:', out_path, '| 共', len(prs.slides.__iter__.__self__._sldIdLst), '页')
